import os
import io
from collections import Counter
from dotenv import load_dotenv
import tiktoken

# Load .env first so that environment variables are available before importing other libs
load_dotenv()

# Force HuggingFace to use a local cache folder to avoid WinError 3 on missing drives (like G:\)
os.environ["HF_HOME"] = os.path.join(os.getcwd(), "hf_cache")

import json
from datetime import datetime
from pathlib import Path
import sys
import re
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
import chromadb
from chromadb.utils import embedding_functions
from openai import OpenAI
import uuid
from werkzeug.utils import secure_filename
from functools import wraps
import time
import threading
import ipaddress
import socket
from urllib.parse import urlparse

jsonl_lock = threading.Lock()

PRIVATE_NETWORKS = [
    ipaddress.ip_network("10.0.0.0/8"),
    ipaddress.ip_network("172.16.0.0/12"),
    ipaddress.ip_network("192.168.0.0/16"),
    ipaddress.ip_network("127.0.0.0/8"),
    ipaddress.ip_network("169.254.0.0/16"),
    ipaddress.ip_network("0.0.0.0/8"),
    ipaddress.ip_network("::1/128"),
    ipaddress.ip_network("fc00::/7"),
]

def is_private_url(url: str) -> bool:
    hostname = urlparse(url).hostname
    if not hostname:
        return True  # reject invalid
    try:
        ip = ipaddress.ip_address(hostname)
    except ValueError:
        # It's a hostname — resolve it
        try:
            ip = ipaddress.ip_address(socket.gethostbyname(hostname))
        except (socket.gaierror, ValueError):
            return True  # reject if can't resolve
    return any(ip in net for net in PRIVATE_NETWORKS)


# Simple Rate Limiter
RATE_LIMITS = {}
def rate_limit(limit_seconds=1.0):
    def decorator(f):
        @wraps(f)
        def wrapped(*args, **kwargs):
            ip = request.remote_addr or "127.0.0.1"
            now = time.time()
            key = (f.__name__, ip)
            last_time = RATE_LIMITS.get(key, 0)
            if now - last_time < limit_seconds:
                return jsonify({"error": "Too many requests"}), 429
            RATE_LIMITS[key] = now
            return f(*args, **kwargs)
        return wrapped
    return decorator

# ---------------------------------------------------------------------------
# Token estimation & context budget
# ---------------------------------------------------------------------------
# Conservative budget for the TOTAL prompt (system + history + user message).
# Leaves headroom for the model's response tokens.
# Small models (e.g. Groq llama-3.1-8b-instant) have limits as low as 6 000 TPM,
# so we aim for ~3 500 input tokens to leave ~2 500 for the response + overhead.
MAX_CONTEXT_TOKENS = int(os.getenv("ULISSE_MAX_CONTEXT_TOKENS", "3500"))

# Cache the encoder
_token_enc = None

def estimate_tokens(text: str) -> int:
    """Precise token counting using tiktoken (cl100k_base)."""
    global _token_enc
    if not text:
        return 0
    if _token_enc is None:
        try:
            # cl100k_base is used by GPT-4 and is a good proxy for LLaMA-3
            _token_enc = tiktoken.get_encoding("cl100k_base")
        except Exception:
            # Fallback to simple estimation if tiktoken fails
            return int(len(text.split()) * 1.33)
    return len(_token_enc.encode(text))

def _parse_dsml_tool_calls(content: str):
    """Parses DeepSeek's leaked <｜｜DSML｜｜tool_calls> format into OpenAI tool call dicts."""
    import re, json, uuid
    parsed = []
    if not content or "<｜｜DSML｜｜tool_calls>" not in content:
        return parsed, content
        
    dsml_match = re.search(r'<｜｜DSML｜｜tool_calls>(.*?)</｜｜DSML｜｜tool_calls>', content, re.DOTALL)
    if not dsml_match:
        return parsed, content
    
    tool_calls_str = dsml_match.group(1)
    invokes = re.finditer(r'<｜｜DSML｜｜invoke name="([^"]+)">(.*?)</｜｜DSML｜｜invoke>', tool_calls_str, re.DOTALL)
    for invoke in invokes:
        fn_name = invoke.group(1)
        params_str = invoke.group(2)
        args_dict = {}
        
        params = re.finditer(r'<｜｜DSML｜｜parameter name="([^"]+)"[^>]*>(.*?)</｜｜DSML｜｜parameter>', params_str, re.DOTALL)
        for param in params:
            args_dict[param.group(1)] = param.group(2).strip()
            
        parsed.append({
            "id": f"call_dsml_{uuid.uuid4().hex[:8]}",
            "type": "function",
            "function": {
                "name": fn_name,
                "arguments": json.dumps(args_dict)
            }
        })
    
    clean_content = content.replace(dsml_match.group(0), "").strip()
    return parsed, clean_content

# --- Optional file-parsing imports ---
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

app = Flask(__name__)
CORS(app)
import torch

# --- MSA Integration ---
from msa_provider import msa_manager
from memory_indexer import indexer

def get_path_from_env(env_var, default_folder):
    env_path = os.getenv(env_var)
    if env_path:
        return Path(env_path)
    script_dir = Path(__file__).resolve().parent
    project_root = script_dir.parent.parent
    return project_root / default_folder

corpus_dir = get_path_from_env("CORPUS_PATH", "corpus")
vectordb_dir = get_path_from_env("VECTORDB_PATH", "vectordb")
chroma_dir = vectordb_dir / "chroma"
frontend_dir = get_path_from_env("FRONTEND_PATH", "webapp/frontend")

corpus_dir.mkdir(parents=True, exist_ok=True)
new_conversations_file = corpus_dir / "new_conversations.jsonl"
sessions_dir = corpus_dir / "sessions"
sessions_dir.mkdir(parents=True, exist_ok=True)
wiki_dir = corpus_dir / "wiki"
wiki_pages_dir = wiki_dir / "pages"
wiki_pages_dir.mkdir(parents=True, exist_ok=True)
wiki_raw_dir = wiki_dir / "raw"
wiki_raw_dir.mkdir(parents=True, exist_ok=True)


# Auto-migration if sessions folder is empty
def check_migration():
    if not any(sessions_dir.iterdir()):
        if not new_conversations_file.exists() or new_conversations_file.stat().st_size == 0:
            return
        print("Sessions folder is empty. Running migration...")
        try:
            from scripts.migrate_to_sessions import migrate
            migrate()
        except ImportError:
            print("Migration script not found in scripts.migrate_to_sessions")
        except Exception as e:
            print(f"Migration error: {e}")

check_migration()

api_key = os.getenv("DEEPSEEK_API_KEY", "")
base_url = os.getenv("DEEPSEEK_BASE_URL", "https://api.deepseek.com")

def get_default_headers(url):
    # Minimal headers to avoid routing issues on some providers
    return None

ai_client = OpenAI(api_key=api_key, base_url=base_url)



chroma_client = None
collection = None
chroma_status = False

def init_chromadb():
    global chroma_client, collection, chroma_status
    try:
        chroma_dir.mkdir(parents=True, exist_ok=True)
        chroma_client = chromadb.PersistentClient(path=str(chroma_dir))
        
        sentence_transformer_ef = embedding_functions.SentenceTransformerEmbeddingFunction(
            model_name="all-MiniLM-L6-v2"
        )
        
        collection = chroma_client.get_or_create_collection(
            name="ulisse_brain",
            embedding_function=sentence_transformer_ef
        )
        chroma_status = True
        print("Connected to ChromaDB 'ulisse_brain' collection.")
    except Exception as e:
        print(f"Error initializing ChromaDB: {e}")
        chroma_status = False

init_chromadb()

script_dir = Path(__file__).resolve().parent
project_root = script_dir.parent.parent
if str(project_root) not in sys.path:
    sys.path.insert(0, str(project_root))

extra_tools = []
extra_tool_handlers = {}

# --- Wiki Tools Implementation ---
def wiki_read_page(args):
    try:
        data = json.loads(args)
        title = data.get("title", "").strip()
        safe_title = re.sub(r'[^\w\s-]', '', title).strip().replace(' ', '_')
        path = (wiki_pages_dir / f"{safe_title}.md").resolve()
        if not str(path).startswith(str(wiki_pages_dir.resolve())):
            return "Error: Access denied (outside of wiki pages directory)."
        if path.exists():
            return f"Content of [[{title}]]:\n\n" + path.read_text(encoding="utf-8")
        return f"The page [[{title}]] does not exist yet."
    except Exception as e:
        return f"Wiki read error: {str(e)}"

def wiki_write_page(args):
    try:
        data = json.loads(args)
        title = data.get("title", "").strip()
        content = data.get("content", "").strip()
        safe_title = re.sub(r'[^\w\s-]', '', title).strip().replace(' ', '_')
        path = wiki_pages_dir / f"{safe_title}.md"
        path.write_text(content, encoding="utf-8")
        return f"Page [[{title}]] updated successfully."
    except Exception as e:
        return f"Wiki write error: {str(e)}"

def wiki_list_pages(args=None):
    try:
        pages = [f.stem.replace('_', ' ') for f in wiki_pages_dir.glob("*.md")]
        return "Pages present in the Wiki:\n" + "\n".join([f"- [[{p}]]" for p in pages])
    except Exception as e:
        return f"Wiki list error: {str(e)}"

def wiki_append_log(args):
    try:
        data = json.loads(args)
        entry = data.get("entry", "").strip()
        log_path = wiki_dir / "log.md"
        with open(log_path, "a", encoding="utf-8") as f:
            f.write(f"\n- [{datetime.now().strftime('%Y-%m-%d %H:%M')}] {entry}")
        return "Log updated."
    except Exception as e:
        return f"Wiki log error: {str(e)}"

def wiki_update_index(args):
    try:
        data = json.loads(args)
        content = data.get("content", "").strip()
        index_path = wiki_dir / "index.md"
        index_path.write_text(content, encoding="utf-8")
        return "Wiki index updated."
    except Exception as e:
        return f"Wiki index error: {str(e)}"

wiki_tools = [
    {
        "type": "function",
        "function": {
            "name": "wiki_read_page",
            "description": "Reads the content of a page from the long-term memory (Wiki).",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "Il titolo della pagina da leggere."}
                },
                "required": ["title"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "wiki_write_page",
            "description": "Creates or updates a page in the long-term memory (Wiki). Use it to store important information, projects, facts, or summaries.",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "Il titolo della pagina."},
                    "content": {"type": "string", "description": "Il contenuto in formato markdown."}
                },
                "required": ["title", "content"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "wiki_list_pages",
            "description": "Lists all pages present in the long-term memory.",
            "parameters": {"type": "object", "properties": {}}
        }
    },
    {
        "type": "function",
        "function": {
            "name": "wiki_append_log",
            "description": "Adds an entry to the activity log (log.md) of the Wiki.",
            "parameters": {
                "type": "object",
                "properties": {
                    "entry": {"type": "string", "description": "Descrizione dell'attività svolta."}
                },
                "required": ["entry"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "wiki_update_index",
            "description": "Updates the central index (index.md) of the Wiki.",
            "parameters": {
                "type": "object",
                "properties": {
                    "content": {"type": "string", "description": "Il nuovo contenuto completo dell'indice."}
                },
                "required": ["content"]
            }
        }
    }
]

extra_tools.extend(wiki_tools)
extra_tool_handlers.update({
    "wiki_read_page": wiki_read_page,
    "wiki_write_page": wiki_write_page,
    "wiki_list_pages": wiki_list_pages,
    "wiki_append_log": wiki_append_log,
    "wiki_update_index": wiki_update_index
})

# --- Native File Tools ---
def native_read_file(args):
    try:
        import json
        data = json.loads(args)
        filepath = data.get("filepath", "")
        # Risolve il path rispetto alla root del progetto
        path = (project_root / filepath).resolve()
        # Verifica di sicurezza basica
        if not str(path).startswith(str(project_root)):
            return "Error: Access denied (outside of project root)."
            
        # Security checks
        ext = path.suffix.lower()
        if ext in ['.key', '.pem', '.pfx', '.p12', '.crt', '.env']:
            return "Error: Access denied (restricted extension)."
        
        restricted_names = ['credentials', 'secrets', 'token', 'id_rsa', 'config.json', '.git-credentials', '.npmrc', 'secrets.json', 'credentials.yml', 'token.txt']
        if path.name.lower() in restricted_names:
            return "Error: Access denied (restricted filename)."
            
        restricted_dirs = ['.git', '.ssh', '.aws', '.config']
        for part in path.parts:
            if part in restricted_dirs:
                return "Error: Access denied (restricted directory)."
                
        if path.exists() and path.is_file():
            if path.stat().st_size > 1024 * 1024:
                return "Error: File too large (max 1MB)."
            return path.read_text(encoding="utf-8", errors="replace")
        return f"File not found: {filepath}"
    except Exception as e:
        return f"Read error: {str(e)}"

def native_list_files(args):
    try:
        import json
        data = json.loads(args)
        directory = data.get("directory", ".")
        path = (project_root / directory).resolve()
        if not str(path).startswith(str(project_root)):
            return "Error: Access denied (outside of project root)."
        if path.exists() and path.is_dir():
            files = []
            for p in path.iterdir():
                rel = p.relative_to(project_root)
                prefix = "📁 " if p.is_dir() else "📄 "
                files.append(f"{prefix}{rel}")
            return "\n".join(files)
        return f"Directory not found: {directory}"
        return f"List error: {str(e)}"
    except Exception as e:
        return f"List error: {str(e)}"

def native_query_memory(args):
    """Cerca nelle conversazioni passate (Short-Term Memory)."""
    try:
        import json
        data = json.loads(args)
        query = data.get("query", "")
        if not query:
            return "Error: No query provided."
            
        if not chroma_status or collection is None:
            return "Error: Memory system (ChromaDB) is not initialized."
            
        results = collection.query(
            query_texts=[query],
            n_results=5
        )
        
        if not results or not results.get("documents") or len(results["documents"][0]) == 0:
            return "No relevant memories found for this query."
            
        docs = results["documents"][0]
        metas = results["metadatas"][0]
        dists = (results.get("distances") or [[]])[0]
        
        output = []
        for doc, meta, dist in zip(docs, metas, dists):
            if dist <= 1.2: # Slightly more permissive than the automatic RAG
                title = meta.get("title", "Untitled")
                date = meta.get("date", "Unknown")
                output.append(f"--- Memory ({title} - {date}) ---\n{doc}")
                
        if not output:
            return "No memories passed the relevance threshold."
            
        return "\n\n".join(output)
    except Exception as e:
        return f"Memory query error: {str(e)}"

native_file_tools = [
    {
        "type": "function",
        "function": {
            "name": "native_read_file",
            "description": "Legge il contenuto di un file nel workspace locale.",
            "parameters": {
                "type": "object",
                "properties": {
                    "filepath": {"type": "string", "description": "Percorso relativo del file da leggere."}
                },
                "required": ["filepath"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "native_list_files",
            "description": "Elenca i file e le cartelle in una directory del workspace.",
            "parameters": {
                "type": "object",
                "properties": {
                    "directory": {"type": "string", "description": "Percorso relativo della directory (default '.')"}
                }
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "native_query_memory",
            "description": "Cerca nelle conversazioni passate e nelle preferenze dell'utente (Short-Term Memory / RAG).",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "La stringa di ricerca per trovare contesti rilevanti."}
                },
                "required": ["query"]
            }
        }
    }
]

extra_tools.extend(native_file_tools)
extra_tool_handlers.update({
    "native_read_file": native_read_file,
    "native_list_files": native_list_files,
    "native_query_memory": native_query_memory
})

# --- Agno Agent Delegation Tool ---
def _make_agno_handler(agno_config):
    """Factory that returns a delegate_to_agno_agent handler with the given
    config bound via closure.  This avoids any shared mutable state
    (threading.local, globals) and is safe for concurrent requests."""
    def delegate_to_agno_agent(args):
        try:
            import json
            data = json.loads(args)
            task = data.get("task", "")
            if not task:
                return "No task provided."

            from webapp.backend.ulisse_agno import get_ulisse_agent

            # Inizializza l'agente con il modello dalla config legata a QUESTA request
            agent = get_ulisse_agent(
                model_id=agno_config.get("model"),
                api_key=agno_config.get("api_key"),
                base_url=agno_config.get("base_url")
            )

            # Prepend retrieved context to the task so Agno has the same knowledge
            full_task = task
            if context_text:
                full_task = f"Context from previous conversations (Short-Term Memory):\n{context_text}\n\nTask: {task}"

            # Esegui l'agente Agno in modo sincrono
            response = agent.run(full_task)
            
            output = ""
            # Capture reasoning if present (passed back to the API/Main Agent)
            if hasattr(response, "reasoning") and response.reasoning:
                output += f"<think>\n{response.reasoning}\n</think>\n\n"
            
            if hasattr(response, "content"):
                output += response.content
                return f"Agno Agent Response:\n{output}"
            return f"Agno Agent Response:\n{str(response)}"
        except Exception as e:
            return f"Error running Agno agent: {str(e)}"
    return delegate_to_agno_agent

agno_tools = [
    {
        "type": "function",
        "function": {
            "name": "delegate_to_agno_agent",
            "description": "Delega un task o una ricerca complessa all'agente Agno. L'agente Agno ha accesso al file system locale tramite il tool Workspace. Usalo per compiti che richiedono di leggere, esplorare o cercare file nel progetto locale.",
            "parameters": {
                "type": "object",
                "properties": {
                    "task": {"type": "string", "description": "Il task dettagliato o la richiesta in linguaggio naturale da delegare all'agente Agno."}
                },
                "required": ["task"]
            }
        }
    }
]

extra_tools.extend(agno_tools)
# NOTE: delegate_to_agno_agent is NOT registered here globally anymore.
# A request-scoped handler is created inside generate() via _make_agno_handler().


# ---------------------------------------------------------------------------
# File text-extraction helpers
# ---------------------------------------------------------------------------

MAX_FILE_SIZE_MB = 15
MAX_EXTRACTED_CHARS = 80_000  # ~20k tokens – safe upper bound for most models

ALLOWED_EXTENSIONS = {
    ".pdf", ".docx", ".pptx", ".xlsx",
    ".txt", ".md", ".csv", ".log", ".json", ".xml", ".yaml", ".yml"
}


def _extract_text_from_file(file_storage) -> str:
    """Extract plain text from a werkzeug FileStorage object.
    Returns the extracted text string (possibly truncated).
    Raises ValueError on unsupported/too-large files.
    """
    filename = secure_filename(file_storage.filename or "upload")
    ext = Path(filename).suffix.lower()

    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(f"Formato file non supportato: '{ext}'. "
                         f"Formati supportati: {', '.join(sorted(ALLOWED_EXTENSIONS))}")

    raw = file_storage.read()
    if len(raw) > MAX_FILE_SIZE_MB * 1024 * 1024:
        raise ValueError(f"File troppo grande (max {MAX_FILE_SIZE_MB} MB).")

    text = ""

    if ext == ".pdf":
        if not HAS_PYMUPDF:
            raise ValueError("PyMuPDF non installato. Esegui: pip install PyMuPDF")
        with fitz.open(stream=raw, filetype="pdf") as doc:
            pages = [page.get_text() for page in doc]
        text = "\n\n".join(pages)

    elif ext == ".docx":
        if not HAS_DOCX:
            raise ValueError("python-docx non installato. Esegui: pip install python-docx")
        doc = DocxDocument(io.BytesIO(raw))
        text = "\n".join(para.text for para in doc.paragraphs)

    elif ext == ".pptx":
        if not HAS_PPTX:
            raise ValueError("python-pptx non installato. Esegui: pip install python-pptx")
        prs = Presentation(io.BytesIO(raw))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        text = "\n".join(parts)

    elif ext == ".xlsx":
        if not HAS_XLSX:
            raise ValueError("openpyxl non installato. Esegui: pip install openpyxl")
        wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        rows = []
        for sheet in wb.worksheets:
            rows.append(f"=== Sheet: {sheet.title} ===")
            for row in sheet.iter_rows(values_only=True):
                row_str = "\t".join(str(c) if c is not None else "" for c in row)
                rows.append(row_str)
        text = "\n".join(rows)

    else:
        # Plain text formats
        text = raw.decode("utf-8", errors="replace")

    # Truncate if needed
    if len(text) > MAX_EXTRACTED_CHARS:
        text = text[:MAX_EXTRACTED_CHARS] + "\n\n[...testo troncato per limiti di dimensione...]"

    return text.strip()


@app.route("/extract_file", methods=["POST"])
@rate_limit(1.0)
def extract_file():
    """Upload a file and return its extracted text as JSON."""
    if "file" not in request.files:
        return jsonify({"error": "Nessun file ricevuto. Usa il campo 'file'."}), 400

    f = request.files["file"]
    if not f or f.filename == "":
        return jsonify({"error": "File non valido."}), 400

    try:
        text = _extract_text_from_file(f)
        return jsonify({
            "filename": secure_filename(f.filename),
            "text": text,
            "char_count": len(text)
        })
    except ValueError as e:
        return jsonify({"error": str(e)}), 422
    except Exception as e:
        print(f"File extraction error: {e}")
        return jsonify({"error": f"Errore durante l'estrazione: {str(e)}"}), 500


# ---------------------------------------------------------------------------

@app.route("/health", methods=["GET"])
def health():
    return jsonify({
        "status": "ok",
        "chromadb": chroma_status
    })

@app.route("/stats", methods=["GET"])
def stats():
    if not chroma_status or collection is None:
        return jsonify({"error": "ChromaDB not connected"}), 500
        
    return jsonify({
        "total_chunks": collection.count(),
        "collection_name": "ulisse_brain"
    })

@app.route("/save_conversation", methods=["POST"])
@rate_limit(1.0)
def save_conversation():
    data = request.json
    if not data:
        return jsonify({"error": "No data provided"}), 400
        
    try:
        with jsonl_lock:
            with open(new_conversations_file, "a", encoding="utf-8") as f:
                f.write(json.dumps(data, ensure_ascii=False) + "\n")
        return jsonify({"status": "success", "message": "Conversation saved."})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

def get_session_path(session_id):
    return sessions_dir / f"{secure_filename(str(session_id))}.json"

# ================================================================
#  MSA MEMORY ENDPOINTS
# ================================================================

@app.route('/msa/index', methods=['POST'])
def msa_index():
    try:
        # Run indexing
        indexer.build_memory_bank()
        return jsonify({"status": "success", "message": "Indexing complete"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/msa/status', methods=['GET'])
def msa_status():
    is_loaded = msa_manager.model is not None
    has_cuda = torch.cuda.is_available()
    return jsonify({
        "is_loaded": is_loaded,
        "device": msa_manager.device,
        "has_cuda": has_cuda,
        "is_loading": msa_manager.is_loading
    })

@app.route("/chat", methods=["POST"])
@rate_limit(1.0)
def chat():
    data = request.json
    user_message = data.get("message", "")
    file_text    = data.get("file_text", "")    # optional extracted file content
    file_name    = data.get("file_name", "")    # optional original filename
    history = data.get("history", []) # Client-side history (optional now, but kept for compatibility)
    session_id = data.get("session_id")

    # Inject extracted file content into the user message
    if file_text:
        label = f"📄 {file_name}" if file_name else "📄 Allegato"
        user_message = (
            f"{user_message}\n\n"
            f"--- {label} ---\n"
            f"{file_text}\n"
            f"--- fine allegato ---"
        ).strip()

    if not user_message:
        return jsonify({"error": "Message is required"}), 400
        
    # Session handling
    is_new_session = False
    if not session_id: # Handles None, "", or missing key
        session_id = str(uuid.uuid4())
        is_new_session = True
        session_data = {
            "id": session_id,
            "title": "New Conversation",
            "created_at": datetime.now().isoformat(),
            "updated_at": datetime.now().isoformat(),
            "messages": [],
            "status": "active"
        }
    else:
        session_path = get_session_path(session_id)
        if session_path.exists():
            with open(session_path, "r", encoding="utf-8") as f:
                session_data = json.load(f)
        else:
            return jsonify({"error": "Session not found"}), 404

    # Use session messages for history instead of what's sent from client if session exists
    chat_history = session_data.get("messages", [])
        
    sources = []
    context_text = ""

    # -------------------------------------------------------------------
    # 1) Load fixed-cost components first to know the baseline budget
    # -------------------------------------------------------------------
    prompt_mode = os.getenv("ULISSE_SYSTEM_PROMPT_MODE", "full")
    prompt_filename = "system_prompt_compact.md" if prompt_mode == "compact" else "system_prompt.md"
    system_prompt_path = corpus_dir / prompt_filename
    
    if not system_prompt_path.exists():
        system_prompt_path = corpus_dir / "system_prompt.md"

    if system_prompt_path.exists():
        base_prompt = system_prompt_path.read_text(encoding="utf-8")
    else:
        base_prompt = "You are Ulisse, an entity that exists, thinks and learns."

    wiki_schema = ""
    schema_path = wiki_dir / "WIKI_SCHEMA.md"
    if schema_path.exists():
        wiki_schema = schema_path.read_text(encoding="utf-8")

    # The boilerplate text that wraps wiki schema and RAG context
    if prompt_mode == "compact":
        system_boilerplate = "[WIKI] You maintain Ulisse's long-term memory. Store important info autonomously. Follow schema below.\n\n"
    else:
        system_boilerplate = (
            "=== MISSION: WIKI MAINTAINER ===\n"
            "You are the keeper of Ulisse's long-term memory. "
            "You must decide autonomously (or upon request) when to store important information in the Wiki. "
            "Strictly follow the schema provided below for Wiki management.\n\n"
        )

    # Fixed token costs
    base_tokens = estimate_tokens(base_prompt)
    schema_tokens = estimate_tokens(wiki_schema)
    boilerplate_tokens = estimate_tokens(system_boilerplate) + 10  # separators
    user_msg_tokens = estimate_tokens(user_message)

    fixed_cost = base_tokens + schema_tokens + boilerplate_tokens + user_msg_tokens
    remaining_budget = max(0, MAX_CONTEXT_TOKENS - fixed_cost)

    print(f"TOKEN BUDGET: total={MAX_CONTEXT_TOKENS}, fixed={fixed_cost} "
          f"(base={base_tokens}, schema={schema_tokens}, user={user_msg_tokens}), "
          f"remaining={remaining_budget}")

    # -------------------------------------------------------------------
    # 2) Split remaining budget: 50% history, 50% RAG context
    # -------------------------------------------------------------------
    history_budget = remaining_budget // 2
    rag_budget = remaining_budget - history_budget

    # -------------------------------------------------------------------
    # 3) Trim chat history — keep most recent messages that fit
    # -------------------------------------------------------------------
    trimmed_history = []
    history_tokens_used = 0
    # Walk backwards (newest first) and collect messages that fit
    for msg in reversed(chat_history):
        role = msg.get("role")
        content = msg.get("content")
        if role in ["user", "assistant"] and content:
            msg_tokens = estimate_tokens(content) + 4  # role overhead
            if history_tokens_used + msg_tokens > history_budget:
                break
            msg_to_add = {"role": role, "content": content}
            if msg.get("reasoning_content"):
                msg_to_add["reasoning_content"] = msg["reasoning_content"]
            trimmed_history.insert(0, msg_to_add)
            history_tokens_used += msg_tokens

    # Any budget the history didn't use goes to RAG
    rag_budget += (history_budget - history_tokens_used)
    
    # If we omitted messages, add a compact summary message
    if len(trimmed_history) < len(chat_history):
        omitted_count = len(chat_history) - len(trimmed_history)
        summary_msg = {
            "role": "system",
            "content": f"[{omitted_count} earlier messages omitted for brevity.]"
        }
        trimmed_history.insert(0, summary_msg)

    print(f"HISTORY: kept {len(trimmed_history)}/{len(chat_history)} messages, "
          f"tokens={history_tokens_used}, RAG budget={rag_budget}")

    # -------------------------------------------------------------------
    # 4) Build RAG context within budget
    # -------------------------------------------------------------------
    if chroma_status and collection is not None:
        try:
            results = collection.query(
                query_texts=[user_message],
                n_results=8
            )

            if results and results.get("documents") and len(results["documents"][0]) > 0:
                retrieved_docs = results["documents"][0]
                retrieved_metas = results["metadatas"][0]
                retrieved_dists = (results.get("distances") or [[]])[0]

                # Filter chunks with reasonable distance (<= 1.0)
                filtered = []
                for doc, meta, dist in zip(retrieved_docs, retrieved_metas, retrieved_dists):
                    if dist <= 1.0:
                        filtered.append((doc, meta))

                print(f"Chunks retrieved: {len(filtered)}/{len(retrieved_docs)} passed threshold")

                context_parts = []
                rag_tokens_used = 0

                for idx, (doc, meta) in enumerate(filtered[:8]):
                    title = meta.get("title", "Untitled")
                    date = meta.get("date", "Unknown")
                    short_date = date[:10] if len(date) > 10 else date

                    if prompt_mode == "compact":
                        header = f"[{idx+1}|{title[:30]}|{short_date}]\n"
                    else:
                        header = f"--- Chunk {idx+1} ({title} - {date}) ---\n"
                    header_tokens = estimate_tokens(header)
                    chunk_tokens = estimate_tokens(doc)

                    if rag_tokens_used + header_tokens >= rag_budget:
                        break

                    if rag_tokens_used + header_tokens + chunk_tokens > rag_budget:
                        # Truncate this chunk to fit
                        allowed_tokens = rag_budget - rag_tokens_used - header_tokens
                        allowed_words = max(0, int(allowed_tokens / 1.33))
                        if allowed_words > 10:
                            truncated = " ".join(doc.split()[:allowed_words]) + "..."
                            context_parts.append(f"{header}{truncated}")
                            sources.append(f"{title} ({date})")
                        break
                    else:
                        context_parts.append(f"{header}{doc}")
                        sources.append(f"{title} ({date})")
                        rag_tokens_used += header_tokens + chunk_tokens

                context_text = "\n\n".join(context_parts)
                sources = list(dict.fromkeys(sources))
                print(f"RAG: included {len(context_parts)} chunks, ~{rag_tokens_used} tokens")
        except Exception as e:
            print(f"ChromaDB search error: {e}")

    # -------------------------------------------------------------------
    # 5) Assemble final messages list
    # -------------------------------------------------------------------
    if prompt_mode == "compact":
        system_prompt = (
            f"{base_prompt}\n\n"
            f"{system_boilerplate}"
            f"{wiki_schema}\n\n"
            f"=== SHORT-TERM MEMORY (STM) ===\n"
            f"The following context was automatically retrieved based on your last message. "
            f"Use it to maintain continuity. If you need more info, use 'native_query_memory'.\n\n"
            f"[MEM]{context_text}\n"
        )
    else:
        system_prompt = (
            f"{base_prompt}\n\n"
            f"{system_boilerplate}"
            f"{wiki_schema}\n\n"
            f"=== RETRIEVED SHORT-TERM MEMORY (STM) ===\n"
            f"Below is the relevant context from previous conversations, retrieved automatically. "
            f"Consider this information as part of your current knowledge. "
            f"If the information is insufficient, use the 'native_query_memory' tool to search further.\n\n"
            f"{context_text}\n"
            f"==========================================\n"
        )

    messages = [{"role": "system", "content": system_prompt}]
    messages.extend(trimmed_history)
    messages.append({"role": "user", "content": user_message})

    total_est = sum(estimate_tokens(m["content"]) for m in messages)
    print(f"FINAL CONTEXT: {len(messages)} messages, ~{total_est} estimated tokens")
    
    # === Provider routing ===
    provider    = data.get("provider", "local")        # local | apikey | ulisse_memo
    req_api_key = data.get("api_key", "")              # only for apikey
    req_base_url= data.get("base_url", "")             # only for apikey
    req_model   = data.get("model", "")                # only for apikey

    # === MSA Provider handling ===
    if provider == "ulisse_memo":
        # Lazy load model on first use
        if msa_manager.model is None:
            print("MSA: First use, loading model...")
            msa_manager.load_model()
            msa_manager.load_memory_bank()
        
        chat_model = "EverMind-AI/MSA-4B"
        # We'll use the generator logic below but bypass OpenAI client
        chat_client = None 
    elif provider == "apikey" and req_api_key:
        # User-supplied API key and base URL
        try:
            target_url = req_base_url or "https://api.openai.com/v1"
            if target_url.endswith('/'):
                target_url = target_url[:-1]
                
            if is_private_url(target_url):
                return jsonify({"error": "Invalid base_url: private or restricted IP."}), 400
                
            if not target_url.startswith("https://") and not request.host.startswith(("127.0.0.1", "localhost")):
                return jsonify({"error": "Insecure connection. API key allows only HTTPS or localhost."}), 400

            user_client = OpenAI(
                api_key=req_api_key,
                base_url=target_url
            )
            chat_client = user_client
            chat_model  = req_model or "gpt-4o"
        except Exception as e:
            return jsonify({"error": f"Invalid API key config: {str(e)}"}), 400
    else:
        # Default: use .env / local model
        # Check if frontend sent local config overrides
        local_base_url = data.get("base_url")
        local_model    = data.get("model")
        
        if local_base_url:
            chat_client = OpenAI(
                api_key="ollama", # dummy for local
                base_url=local_base_url
            )
            chat_model = local_model or "llama3"
        else:
            chat_client = ai_client
            chat_model  = "deepseek-chat"

    # === LLM call (local / apikey paths) ===
    print(f"DEBUG: Provider={provider}, Model={chat_model}, BaseURL={chat_client.base_url}")
    print(f"DEBUG: Headers={chat_client.default_headers}")

    from flask import Response
    import copy
    
    def generate():
        nonlocal session_data, messages, sources
        
        # Build request-scoped tool handlers dict.
        # The Agno delegate gets the current request's model config via closure
        # — no shared mutable state, fully thread-safe.
        request_agno_config = {
            "model": chat_model,
            "api_key": chat_client.api_key if chat_client else "msa",
            "base_url": str(chat_client.base_url) if chat_client else "local"
        }
        request_tool_handlers = dict(extra_tool_handlers)
        request_tool_handlers["delegate_to_agno_agent"] = _make_agno_handler(request_agno_config)
        
        # Invia l'inizio della sessione
        yield f"data: {json.dumps({'event': 'session', 'session_id': session_id, 'session_title': session_data.get('title', 'Nuova Chat'), 'sources': sources})}\n\n"
        
        full_assistant_response = ""
        full_assistant_reasoning = ""
        
        try:
            # Pass 1: Tool Routing
            yield f"data: {json.dumps({'event': 'think', 'message': 'Analisi in corso...'})}\n\n"
            
            # Check if we should enable reasoning for OpenRouter/DeepSeek-R1
            extra_body = {}
            if "r1" in chat_model.lower() and "openrouter" in str(chat_client.base_url).lower():
                extra_body["reasoning"] = {"enabled": True}

            # Tools are kept enabled for OpenRouter as many free models support them
            current_tools = extra_tools if extra_tools else None
            is_openrouter = "openrouter.ai" in str(chat_client.base_url)

            print(f"DEBUG: extra_body={extra_body}")

            def _llm_call(msgs, tools, stream=False, max_tok=4000, retry_count=0):
                """Wrapper that handles tool_use_failed by retrying without tools."""
                if retry_count > 2:
                    raise Exception("Too many retries in LLM call")

                kwargs = dict(model=chat_model, messages=msgs, stream=stream)
                if not is_openrouter:
                    kwargs["temperature"] = 0.7
                    kwargs["max_tokens"] = max_tok
                if tools:
                    kwargs["tools"] = tools
                if extra_body:
                    kwargs["extra_body"] = extra_body

                try:
                    if provider == "ulisse_memo":
                        # Convert messages to MSA format if needed (usually just standard list of dicts)
                        # We use the native transformer/tokenizer generate
                        res = msa_manager.generate(msgs, tools=tools)
                        
                        # Mock the OpenAI response object structure
                        class MockChoice:
                            def __init__(self, content, tool_calls):
                                self.message = type('obj', (object,), {'content': content, 'tool_calls': tool_calls, 'model_dump': lambda: {'content': content, 'tool_calls': tool_calls}})
                        
                        class MockResponse:
                            def __init__(self, content, tool_calls):
                                self.choices = [MockChoice(content, tool_calls)]
                        
                        return MockResponse(res.get("content"), res.get("tool_calls")), tools
                    else:
                        return chat_client.chat.completions.create(**kwargs), tools
                except Exception as e:
                    err_str = str(e)
                    # Model can't handle tool calling → retry without tools
                    if "tool_use_failed" in err_str or "tool_calls" in err_str:
                        print(f"Tool calling failed ({e}), retrying WITHOUT tools")
                        kwargs.pop("tools", None)
                        return _llm_call(msgs, None, stream, max_tok, retry_count + 1)
                    # Token limit → retry with lower max_tokens
                    if "413" in err_str or "rate_limit_exceeded" in err_str or "tokens per minute" in err_str.lower():
                        if max_tok > 1500:
                            print(f"Token limit API Error: {e}, retrying with max_tokens=1500")
                            return _llm_call(msgs, tools, stream, 1500, retry_count + 1)
                    raise

            response, effective_tools = _llm_call(messages, current_tools)
            assistant_message = response.choices[0].message
            
            # Capture and yield reasoning content if present during tool calling phase
            msg_reasoning = getattr(assistant_message, 'reasoning_content', None) or (assistant_message.model_dump().get('reasoning_content') if hasattr(assistant_message, 'model_dump') else None)
            if msg_reasoning:
                full_assistant_reasoning += msg_reasoning
                yield f"data: {json.dumps({'event': 'think', 'message': msg_reasoning, 'append': True})}\n\n"

            # Detect models that output tool calls as plain text instead of
            # using the structured tool_calls API (common with small models).
            # Pattern: <function=name>{...}</function> in the content text.
            _fn_leak_re = re.compile(r'<function=\w+>.*?</function>', re.DOTALL)
            if (assistant_message.content
                    and _fn_leak_re.search(assistant_message.content)
                    and not assistant_message.tool_calls):
                print("DETECTED: model leaked tool calls as text, retrying WITHOUT tools")
                effective_tools = None
                response, _ = _llm_call(messages, None)
                assistant_message = response.choices[0].message
            
            # --- DSML Interceptor Pass ---
            dsml_parsed_tools, clean_content = _parse_dsml_tool_calls(assistant_message.content) if assistant_message.content else ([], assistant_message.content)
            
            has_tools = bool(assistant_message.tool_calls) or bool(dsml_parsed_tools)

            while has_tools:
                # Convert to dict and ensure content is handled for API compatibility
                msg_dict = assistant_message.model_dump(exclude_none=True)
                
                # Clean leak from content before adding to history
                if msg_dict.get("content"):
                    msg_dict["content"] = _fn_leak_re.sub('', msg_dict["content"]).strip()
                    if not msg_dict["content"]:
                        msg_dict["content"] = None
                
                # Remove fields not supported in requests by some providers
                for field in ["annotations", "audio", "refusal"]:
                    msg_dict.pop(field, None)
                
                # IMPORTANT: DeepSeek-R1/OpenRouter REQUIRE reasoning_content to be passed back
                if msg_reasoning:
                    msg_dict["reasoning_content"] = msg_reasoning
                    
                # Apply DSML parsing if present
                if dsml_parsed_tools:
                    msg_dict["content"] = clean_content if clean_content else None
                    msg_dict["tool_calls"] = dsml_parsed_tools

                messages.append(msg_dict)
                
                tool_calls_list = msg_dict.get("tool_calls", [])
                for tc in tool_calls_list:
                    # Handle both native ChatCompletionMessageToolCall dict dumps and DSML dicts
                    if isinstance(tc, dict):
                        tc_id = tc["id"]
                        fn_name = tc["function"]["name"]
                        fn_args = tc["function"]["arguments"]
                    else:
                        tc_id = tc.id
                        fn_name = tc.function.name
                        fn_args = tc.function.arguments

                    yield f"data: {json.dumps({'event': 'tool', 'message': f'Azione: {fn_name}'})}\n\n"
                    handler = request_tool_handlers.get(fn_name)
                    
                    if handler:
                        try:
                            result = handler(fn_args)
                        except Exception as e:
                            result = f"Error executing tool {fn_name}: {str(e)}"
                    else:
                        result = f"Error: Tool '{fn_name}' not found."
                    
                    messages.append({
                        "role": "tool",
                        "tool_call_id": tc_id,
                        "name": fn_name,
                        "content": str(result)
                    })
                
                yield f"data: {json.dumps({'event': 'think', 'message': 'Elaborazione...'})}\n\n"
                response, effective_tools = _llm_call(messages, effective_tools)
                assistant_message = response.choices[0].message
                
                # Capture and yield reasoning content if present during tool calling loop
                msg_reasoning = getattr(assistant_message, 'reasoning_content', None) or (assistant_message.model_dump().get('reasoning_content') if hasattr(assistant_message, 'model_dump') else None)
                if msg_reasoning:
                    full_assistant_reasoning += msg_reasoning
                    yield f"data: {json.dumps({'event': 'think', 'message': msg_reasoning, 'append': True})}\n\n"
                
                # Re-check for next loop iteration
                dsml_parsed_tools, clean_content = _parse_dsml_tool_calls(assistant_message.content) if assistant_message.content else ([], assistant_message.content)
                has_tools = bool(assistant_message.tool_calls) or bool(dsml_parsed_tools)
            
            # Pass 2: Generazione del testo finale (Stream)
            yield f"data: {json.dumps({'event': 'think', 'message': 'Rispondo...'})}\n\n"
            
            stream_response, _ = _llm_call(messages, None, stream=True)
            
            for chunk in stream_response:
                if not chunk.choices:
                    continue
                delta = chunk.choices[0].delta
                
                reasoning = getattr(delta, 'reasoning_content', None) or (delta.model_dump().get('reasoning_content') if hasattr(delta, 'model_dump') else None)
                if reasoning:
                    full_assistant_reasoning += reasoning
                    yield f"data: {json.dumps({'event': 'think', 'message': reasoning, 'append': True})}\n\n"

                if delta.content:
                    full_assistant_response += delta.content
                    yield f"data: {json.dumps({'event': 'message', 'text': delta.content})}\n\n"

            # Post-process: strip any leaked function-call text from saved response
            full_assistant_response = _fn_leak_re.sub('', full_assistant_response).strip()
                    
        except Exception as e:
            print(f"LLM API error: {e}")
            yield f"data: {json.dumps({'event': 'error', 'message': f'Failed to generate response: {str(e)}'})}\n\n"
            if not full_assistant_response:
                full_assistant_response = "Errore di generazione."

        # --- Post-processing e Salvataggio ---
        timestamp = datetime.now().isoformat()
        
        # Auto-generate title from first message
        if len(chat_history) == 0:
            words = user_message.split()
            if len(words) >= 3:
                session_data["title"] = " ".join(w.capitalize() for w in words[:6])
            else:
                session_data["title"] = f"Chat {timestamp[:16]}"
        
        # AI Title refinement after 3 messages
        if len(chat_history) == 4: # 2 user + 2 assistant already, this is the 3rd pair
            try:
                title_prompt = [
                    {"role": "system", "content": "Generate a synthetic title (max 6 words) for this conversation between Ulisse and Toni. Respond ONLY with the title."},
                    {"role": "user", "content": f"Initial messages:\n" + "\n".join([f"{m['role']}: {m['content'][:100]}" for m in chat_history[:4]]) + f"\nuser: {user_message}"}
                ]
                t_resp = chat_client.chat.completions.create(
                    model=chat_model,
                    messages=title_prompt,
                    max_tokens=20
                )
                new_title = t_resp.choices[0].message.content.strip().strip('"')
                if new_title:
                    session_data["title"] = new_title
            except Exception as e:
                print(f"Title generation error: {e}")

        # Append to session
        session_data["messages"].append({
            "role": "user",
            "content": user_message,
            "timestamp": timestamp
        })
        assistant_msg = {
            "role": "assistant",
            "content": full_assistant_response,
            "timestamp": timestamp,
            "sources": sources
        }
        if full_assistant_reasoning:
            assistant_msg["reasoning_content"] = full_assistant_reasoning
            
        session_data["messages"].append(assistant_msg)
        session_data["updated_at"] = timestamp
        
        # Save session
        try:
            with open(get_session_path(session_id), "w", encoding="utf-8") as f:
                json.dump(session_data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"Error saving session: {e}")

        # Save to legacy jsonl for RAG
        chunk_id = f"chunk_{session_id}_{len(session_data['messages']) // 2}"
        exchange = {
            "session_id": session_id,
            "title": session_data["title"],
            "timestamp": timestamp,
            "user_message": user_message,
            "assistant_response": full_assistant_response,
            "sources": sources,
            "chunk_id": chunk_id
        }
        try:
            with jsonl_lock:
                with open(new_conversations_file, "a", encoding="utf-8") as f:
                    f.write(json.dumps(exchange, ensure_ascii=False) + "\n")
                
            # Real-time embedding into ChromaDB (Short-term memory)
            if chroma_status and collection is not None:
                document = f"User: {user_message}\nAssistant: {full_assistant_response}"
                collection.upsert(
                    documents=[document],
                    metadatas=[{
                        "session_id": str(session_id),
                        "title": str(session_data["title"]) if session_data["title"] else "Untitled",
                        "date": str(timestamp),
                        "type": "conversation"
                    }],
                    ids=[chunk_id]
                )
                print(f"Real-time STM update: Added chunk {chunk_id} to ChromaDB.")
                
        except Exception as e:
            print(f"Error saving legacy exchange or updating STM: {e}")

        # Segnale di fine stream
        yield f"data: {json.dumps({'event': 'done', 'session_title': session_data['title'], 'session_id': session_id})}\n\n"

    return Response(generate(), mimetype="text/event-stream")

@app.route("/sessions", methods=["GET"])
def get_sessions():
    sessions = []
    for f in sessions_dir.glob("*.json"):
        try:
            with open(f, "r", encoding="utf-8") as file:
                data = json.load(file)
                sessions.append({
                    "id": data["id"],
                    "title": data["title"],
                    "created_at": data.get("created_at"),
                    "updated_at": data.get("updated_at"),
                    "message_count": len(data.get("messages", [])) // 2
                })
        except:
            continue
    # Sort by updated_at descending
    sessions.sort(key=lambda x: x.get("updated_at", ""), reverse=True)
    return jsonify(sessions)

@app.route("/sessions/<session_id>", methods=["GET", "DELETE", "PATCH"])
def handle_session(session_id):
    path = get_session_path(session_id)
    if not path.exists():
        return jsonify({"error": "Session not found"}), 404
        
    if request.method == "GET":
        with open(path, "r", encoding="utf-8") as f:
            return jsonify(json.load(f))
            
    elif request.method == "DELETE":
        path.unlink()
        return jsonify({"status": "deleted"})
        
    elif request.method == "PATCH":
        data = request.json
        with open(path, "r", encoding="utf-8") as f:
            session_data = json.load(f)
        if "title" in data:
            session_data["title"] = data["title"]
        with open(path, "w", encoding="utf-8") as f:
            json.dump(session_data, f, ensure_ascii=False, indent=2)
        return jsonify({"status": "updated", "title": session_data["title"]})

@app.route("/memory/stats", methods=["GET"])
def get_memory_stats():
    sessions = []
    total_messages = 0
    all_titles = []
    
    for f in sessions_dir.glob("*.json"):
        try:
            with open(f, "r", encoding="utf-8") as file:
                data = json.load(file)
                sessions.append(data)
                total_messages += len(data.get("messages", []))
                all_titles.append(data.get("title", ""))
        except:
            continue
            
    if not sessions:
        return jsonify({
            "total_sessions": 0,
            "total_chunks": collection.count() if collection else 0,
            "total_messages": 0,
            "oldest_session": None,
            "newest_session": None,
            "most_active_topics": []
        })

    # Sort sessions by creation date
    sessions_sorted = sorted(sessions, key=lambda x: x.get("created_at", ""))
    
    # Extract keywords for topics
    words = []
    for title in all_titles:
        # Simple tokenization: lower case, alphanumeric only, length > 3
        found = re.findall(r'\b[a-z]{4,}\b', title.lower())
        words.extend(found)
    
    # Filter out common Italian stop words if necessary, but for now just top 10
    top_topics = [item[0] for item in Counter(words).most_common(10)]
    
    return jsonify({
        "total_sessions": len(sessions),
        "total_chunks": collection.count() if collection else 0,
        "total_messages": total_messages,
        "oldest_session": sessions_sorted[0].get("created_at"),
        "newest_session": sessions_sorted[-1].get("created_at"),
        "most_active_topics": top_topics
    })

@app.route("/memory/nodes", methods=["GET"])
def get_memory_nodes():
    nodes = []
    edges = []
    
    # Load Wiki Pages instead of sessions for the Memory View
    wiki_pages = list(wiki_pages_dir.glob("*.md"))
    
    for f in wiki_pages:
        try:
            title = f.stem.replace('_', ' ')
            content = f.read_text(encoding="utf-8")
            
            # Extract basic stats
            word_count = len(content.split())
            
            node = {
                "id": f.stem,
                "title": title,
                "content": content,
                "type": "wiki",
                "weight": max(5, min(word_count // 10, 20)),
                "created_at": datetime.fromtimestamp(f.stat().st_ctime).isoformat(),
                "updated_at": datetime.fromtimestamp(f.stat().st_mtime).isoformat()
            }
            nodes.append(node)
        except:
            continue
            
    # Generate Edges based on Wiki Links [[Page Name]]
    page_ids = [n["id"] for n in nodes]
    for n in nodes:
        content = n["content"]
        # Find all [[...]] links
        links = re.findall(r'\[\[([^\]]+)\]\]', content)
        for link in links:
            target_id = re.sub(r'[^\w\s-]', '', link).strip().replace(' ', '_')
            if target_id in page_ids and target_id != n["id"]:
                edges.append({
                    "source": n["id"],
                    "target": target_id,
                    "weight": 1
                })
                
    return jsonify({"nodes": nodes, "edges": edges})
    
@app.route("/memory/nodes/<node_id>", methods=["POST", "DELETE", "PATCH"])
def handle_memory_node(node_id):
    # node_id is the filename stem (e.g., 'Project_Alpha')
    file_path = wiki_pages_dir / f"{node_id}.md"
    
    if request.method == "DELETE":
        if file_path.exists():
            file_path.unlink()
            return jsonify({"status": "deleted"})
        return jsonify({"error": "Node not found"}), 404
        
    elif request.method == "POST":
        data = request.json
        new_content = data.get("content")
        if new_content is not None:
            file_path.write_text(new_content, encoding="utf-8")
            return jsonify({"status": "updated"})
        return jsonify({"error": "No content provided"}), 400

    elif request.method == "PATCH":
        data = request.json
        new_title = data.get("title")
        if new_title:
            # Generate new safe filename
            new_id = re.sub(r'[^\w\s-]', '', new_title).strip().replace(' ', '_')
            new_path = wiki_pages_dir / f"{new_id}.md"
            if file_path.exists():
                file_path.rename(new_path)
                return jsonify({"status": "renamed", "new_id": new_id})
            return jsonify({"error": "Node not found"}), 404
        return jsonify({"error": "No title provided"}), 400




@app.route("/sessions/<session_id>/related", methods=["GET"])
def get_related_sessions(session_id):
    path = get_session_path(session_id)
    if not path.exists() or not chroma_status or collection is None:
        return jsonify([])
        
    with open(path, "r", encoding="utf-8") as f:
        session_data = json.load(f)
    
    # Get last user message to find related
    msgs = session_data.get("messages", [])
    user_msgs = [m["content"] for m in msgs if m["role"] == "user"]
    if not user_msgs:
        return jsonify([])
    
    query = user_msgs[-1]
    try:
        results = collection.query(query_texts=[query], n_results=10)
        metas = results.get("metadatas", [[]])[0]
        
        related_ids = set()
        for m in metas:
            sid = m.get("session_id")
            if sid and sid != session_id:
                related_ids.add(sid)
        
        # Build session list
        related = []
        for sid in list(related_ids)[:5]:
            p = get_session_path(sid)
            if p.exists():
                with open(p, "r", encoding="utf-8") as f:
                    d = json.load(f)
                    related.append({"id": d["id"], "title": d["title"]})
        return jsonify(related)
    except:
        return jsonify([])

# Synthetic memory routes removed as per user request


@app.route("/graph", methods=["GET"])
def get_graph():
    try:
        sessions_list = []
        for f in sessions_dir.glob("*.json"):
            try:
                with open(f, "r", encoding="utf-8") as file:
                    sessions_list.append(json.load(file))
            except:
                continue
        
        nodes = []
        for s in sessions_list:
            sid = s["id"]
            title = s["title"]
            msgs = s.get("messages", [])
            msg_count = len(msgs) // 2
            
            nodes.append({
                "id": sid,
                "label": title,
                "type": "session",
                "weight": max(5, min(msg_count * 2, 20)),
                "created_at": s.get("created_at"),
                "updated_at": s.get("updated_at")
            })
            
        edges = []
        if chroma_status and collection is not None and len(sessions_list) > 1:
            # Generate edges based on RAG similarity
            for i, s1 in enumerate(sessions_list):
                sid1 = s1["id"]
                title1 = s1.get("title", "")
                
                try:
                    # Query ChromaDB for sessions similar to this one's title or first message
                    query_text = title1
                    if s1.get("messages"):
                        query_text += " " + s1["messages"][0]["content"]
                    
                    results = collection.query(
                        query_texts=[query_text],
                        n_results=5,
                        where={"session_id": {"$ne": sid1}}
                    )
                    
                    if results and results.get("metadatas"):
                        metas = results["metadatas"][0]
                        dists = results.get("distances", [[]])[0]
                        
                        for meta, dist in zip(metas, dists):
                            sid2 = meta.get("session_id")
                            if sid2:
                                # Connection strength based on distance (closer to 0 is stronger)
                                strength = max(0.1, 1.0 - dist)
                                edges.append({
                                    "from": sid1,
                                    "to": sid2,
                                    "strength": strength
                                })
                except Exception as e:
                    print(f"Error calculating graph edges for {sid1}: {e}")
        
        # Deduplicate edges (A->B and B->A might exist)
        unique_edges = {}
        for e in edges:
            pair = tuple(sorted([e["from"], e["to"]]))
            if pair not in unique_edges or e["strength"] > unique_edges[pair]["strength"]:
                unique_edges[pair] = e
        
        return jsonify({
            "nodes": nodes,
            "edges": list(unique_edges.values())
        })
    except Exception as e:
        print(f"Graph generation error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/", methods=["GET"])
def index():
    return send_from_directory(str(frontend_dir), "index.html")

@app.route("/<path:filename>", methods=["GET"])
def serve_static(filename):
    return send_from_directory(str(frontend_dir), filename)

if __name__ == "__main__":
    # Disabilitiamo il reloader automatico per evitare conflitti con i processi in background (Playwright/Agno) su Windows
    app.run(host="0.0.0.0", port=5000, debug=True, use_reloader=False)
